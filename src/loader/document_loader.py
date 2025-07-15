"""
Document loader with multi-format parsing and security sanitization.
Production-ready implementation for Erguvan AI Content Generator.
"""

import os
import re
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from datetime import datetime
import hashlib

import pdfplumber
from docx import Document
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup
import bleach

from src.config import config_manager


@dataclass
class DocumentMetadata:
    """Document metadata structure."""
    file_path: str
    file_name: str
    file_size: int
    file_type: str
    created_at: datetime
    modified_at: datetime
    language: str
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    document_id: str = ""
    
    def __post_init__(self):
        """Generate document ID after initialization."""
        if not self.document_id:
            content = f"{self.file_path}{self.file_size}{self.modified_at}"
            self.document_id = hashlib.md5(content.encode()).hexdigest()


@dataclass
class DocumentChunk:
    """Document chunk with metadata for vector storage."""
    content: str
    metadata: DocumentMetadata
    chunk_index: int
    heading: Optional[str] = None
    page_number: Optional[int] = None
    start_char: int = 0
    end_char: int = 0
    
    def __post_init__(self):
        """Calculate word count for the chunk."""
        self.word_count = len(self.content.split())


class DocumentSanitizer:
    """Security sanitization for document content."""
    
    # Allowed HTML tags for cleaning
    ALLOWED_TAGS = ['p', 'br', 'strong', 'em', 'u', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'ul', 'ol', 'li']
    
    # Patterns for removing boilerplate content
    BOILERPLATE_PATTERNS = [
        r'©\s*\d{4}.*?all rights reserved',
        r'confidential.*?distribution',
        r'this document.*?proprietary',
        r'legal notice.*?disclaimer',
        r'copyright.*?\d{4}',
        r'page \d+ of \d+',
        r'printed on.*?\d{4}',
        r'document version.*?\d+\.\d+',
    ]
    
    # Patterns for sensitive information
    SENSITIVE_PATTERNS = [
        r'\b\d{4}[-\s]?\d{4}[-\s]?\d{4}[-\s]?\d{4}\b',  # Credit card numbers
        r'\b\d{3}-\d{2}-\d{4}\b',  # SSN
        r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',  # Email addresses
        r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b',  # Phone numbers
    ]
    
    @classmethod
    def sanitize_content(cls, content: str) -> str:
        """Sanitize document content by removing boilerplate and sensitive info."""
        # Remove HTML tags and normalize whitespace
        clean_content = bleach.clean(content, tags=cls.ALLOWED_TAGS, strip=True)
        
        # Remove boilerplate patterns
        for pattern in cls.BOILERPLATE_PATTERNS:
            clean_content = re.sub(pattern, '', clean_content, flags=re.IGNORECASE)
        
        # Redact sensitive information
        for pattern in cls.SENSITIVE_PATTERNS:
            clean_content = re.sub(pattern, '[REDACTED]', clean_content)
        
        # Normalize whitespace
        clean_content = re.sub(r'\s+', ' ', clean_content)
        clean_content = re.sub(r'\n\s*\n', '\n', clean_content)
        
        return clean_content.strip()
    
    @classmethod
    def validate_file_security(cls, file_path: Path) -> bool:
        """Validate file security and size constraints."""
        config = config_manager.load_models_config()
        
        # Check file size
        max_size = int(os.getenv('MAX_FILE_SIZE_MB', '10')) * 1024 * 1024
        if file_path.stat().st_size > max_size:
            raise ValueError(f"File size exceeds maximum of {max_size/1024/1024}MB")
        
        # Check file extension
        allowed_extensions = os.getenv('ALLOWED_FILE_TYPES', '.pdf,.docx,.md,.pptx').split(',')
        if file_path.suffix not in allowed_extensions:
            raise ValueError(f"File type {file_path.suffix} not allowed")
        
        return True


class DocumentLoader:
    """Production-ready document loader with security and parsing capabilities."""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.sanitizer = DocumentSanitizer()
        
    def load_document(self, file_path: str) -> Tuple[str, DocumentMetadata]:
        """
        Load a document from file path and return content with metadata.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Tuple of (content, metadata)
            
        Raises:
            ValueError: If file is invalid or insecure
            FileNotFoundError: If file doesn't exist
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        # Security validation
        self.sanitizer.validate_file_security(path)
        
        # Get file stats
        stat = path.stat()
        
        # Parse content based on file type
        content = self._parse_content(path)
        
        # Sanitize content
        content = self.sanitizer.sanitize_content(content)
        
        # Detect language
        language = self._detect_language(content)
        
        # Create metadata
        metadata = DocumentMetadata(
            file_path=str(path),
            file_name=path.name,
            file_size=stat.st_size,
            file_type=path.suffix,
            created_at=datetime.fromtimestamp(stat.st_ctime),
            modified_at=datetime.fromtimestamp(stat.st_mtime),
            language=language,
            word_count=len(content.split())
        )
        
        self.logger.info(f"Loaded document: {path.name} ({metadata.word_count} words)")
        
        return content, metadata
    
    def chunk_document(self, content: str, metadata: DocumentMetadata, 
                      chunk_size: int = 400, overlap: int = 50) -> List[DocumentChunk]:
        """
        Split document content into chunks with overlap for vector storage.
        
        Args:
            content: Document content
            metadata: Document metadata
            chunk_size: Maximum tokens per chunk
            overlap: Overlap tokens between chunks
            
        Returns:
            List of DocumentChunk objects
        """
        chunks = []
        
        # Split by headings first (if available)
        sections = self._split_by_headings(content)
        
        chunk_index = 0
        for section_heading, section_content in sections:
            # Split section into smaller chunks if needed
            words = section_content.split()
            
            start_idx = 0
            while start_idx < len(words):
                end_idx = min(start_idx + chunk_size, len(words))
                chunk_words = words[start_idx:end_idx]
                chunk_content = ' '.join(chunk_words)
                
                # Create chunk
                chunk = DocumentChunk(
                    content=chunk_content,
                    metadata=metadata,
                    chunk_index=chunk_index,
                    heading=section_heading,
                    start_char=start_idx,
                    end_char=end_idx
                )
                chunks.append(chunk)
                
                chunk_index += 1
                start_idx = end_idx - overlap
        
        self.logger.info(f"Created {len(chunks)} chunks from document: {metadata.file_name}")
        
        return chunks
    
    def _parse_content(self, file_path: Path) -> str:
        """Parse content based on file type."""
        extension = file_path.suffix.lower()
        
        if extension == '.pdf':
            return self._parse_pdf(file_path)
        elif extension == '.docx':
            return self._parse_docx(file_path)
        elif extension == '.pptx':
            return self._parse_pptx(file_path)
        elif extension == '.md':
            return self._parse_markdown(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")
    
    def _parse_pdf(self, file_path: Path) -> str:
        """Parse PDF content using pdfplumber."""
        content = []
        
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
        
        return '\n'.join(content)
    
    def _parse_docx(self, file_path: Path) -> str:
        """Parse DOCX content using python-docx."""
        doc = Document(file_path)
        content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        
        return '\n'.join(content)
    
    def _parse_pptx(self, file_path: Path) -> str:
        """Parse PPTX content using python-pptx."""
        prs = Presentation(file_path)
        content = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    content.append(shape.text)
        
        return '\n'.join(content)
    
    def _parse_markdown(self, file_path: Path) -> str:
        """Parse Markdown content."""
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        # Convert to HTML and extract text
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, 'html.parser')
        
        return soup.get_text()
    
    def _split_by_headings(self, content: str) -> List[Tuple[Optional[str], str]]:
        """Split content by headings to preserve structure."""
        # Simple heading detection patterns
        heading_patterns = [
            r'^#{1,6}\s+(.+)$',  # Markdown headings
            r'^([A-Z][A-Z\s]+)$',  # ALL CAPS headings
            r'^\d+\.\s+(.+)$',  # Numbered headings
        ]
        
        sections = []
        current_heading = None
        current_content = []
        
        for line in content.split('\n'):
            line = line.strip()
            if not line:
                continue
            
            # Check if line is a heading
            is_heading = False
            for pattern in heading_patterns:
                match = re.match(pattern, line)
                if match:
                    # Save previous section
                    if current_content:
                        sections.append((current_heading, '\n'.join(current_content)))
                    
                    # Start new section
                    current_heading = match.group(1) if match.groups() else line
                    current_content = []
                    is_heading = True
                    break
            
            if not is_heading:
                current_content.append(line)
        
        # Add final section
        if current_content:
            sections.append((current_heading, '\n'.join(current_content)))
        
        return sections if sections else [(None, content)]
    
    def _detect_language(self, content: str) -> str:
        """Simple language detection based on common words."""
        # Turkish common words
        turkish_words = ['ve', 'bir', 'bu', 'ile', 'için', 'olan', 'olarak', 'çok', 'daha', 'karbon']
        
        # English common words
        english_words = ['and', 'the', 'for', 'with', 'this', 'that', 'are', 'carbon', 'climate', 'policy']
        
        words = content.lower().split()
        
        turkish_count = sum(1 for word in words if word in turkish_words)
        english_count = sum(1 for word in words if word in english_words)
        
        return 'tr' if turkish_count > english_count else 'en'


class DocumentBatch:
    """Batch processing for multiple documents."""
    
    def __init__(self, loader: DocumentLoader):
        self.loader = loader
        self.logger = logging.getLogger(__name__)
    
    def load_directory(self, directory_path: str) -> List[Tuple[str, DocumentMetadata]]:
        """Load all supported documents from a directory."""
        directory = Path(directory_path)
        
        if not directory.exists():
            raise FileNotFoundError(f"Directory not found: {directory_path}")
        
        documents = []
        supported_extensions = ['.pdf', '.docx', '.pptx', '.md']
        
        for file_path in directory.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                try:
                    content, metadata = self.loader.load_document(str(file_path))
                    documents.append((content, metadata))
                except Exception as e:
                    self.logger.error(f"Failed to load {file_path}: {e}")
        
        self.logger.info(f"Loaded {len(documents)} documents from {directory_path}")
        
        return documents
    
    def process_batch(self, file_paths: List[str]) -> List[Tuple[str, DocumentMetadata]]:
        """Process a batch of document files."""
        documents = []
        
        for file_path in file_paths:
            try:
                content, metadata = self.loader.load_document(file_path)
                documents.append((content, metadata))
            except Exception as e:
                self.logger.error(f"Failed to load {file_path}: {e}")
        
        return documents